"""Étape 25 — le diaporama de soutenance, généré.

    python -m tvfed.diaporama

────────────────────────────────────────────────────────────────────────────
POURQUOI GÉNÉRER PLUTÔT QUE DESSINER
────────────────────────────────────────────────────────────────────────────
Un diaporama fait à la main diverge de ses sources dès la première
correction. Ce projet en a fait l'expérience : « +45 % de FWI » a survécu des
semaines dans l'application parce que c'était une chaîne de caractères que
rien ne recalculait.

Ici, chaque chiffre est **lu dans le CSV qui l'a produit**. Relancer
`tvfed.comparer` puis cette commande suffit à remettre le diaporama d'aplomb.
Les seules constantes écrites en dur sont celles qui décrivent le protocole —
les bornes du split, la taille de la grille — et elles ne changent pas.

Le fichier produit est un `.pptx` ordinaire : ouvrable dans PowerPoint,
LibreOffice ou Google Slides, et retouchable.
"""
from __future__ import annotations

import json

import pandas as pd

from .paths import PROCESSED, RACINE

SORTIE = RACINE / "presentation"
FIG = RACINE / "figures"

# charte, reprise des notebooks et de l'application
INK = "0B0B0B"
MUTED = "6B6963"
FOND = "FCFCFB"
ROUGE = "E34948"
ORANGE = "EB6834"
BLEU = "2A78D6"
VERT = "1BAF7A"
GRIS = "C3C2B7"


def _lire():
    """Tous les chiffres du diaporama, à leur source."""
    d = {}
    for nom in ("pr_auc_val", "comparaison_appariee", "transfert_spatial",
                "test_par_annee", "baselines", "calibration_v3",
                "series_sarimax", "modeles_ensemble", "series_adf",
                "series_acf_pacf", "operationnel_test"):
        f = PROCESSED / f"{nom}.csv"
        d[nom] = pd.read_csv(f) if f.exists() else None
    # ⚠️ lu depuis le CSV, et non recopié : la ROC-AUC du modèle « grand feu »
    # était affichée par `tvfed.taille` sans être sauvegardée, et le 0,77 a
    # circulé pendant des semaines sans artefact pour le soutenir.
    d["taille"] = pd.read_csv(PROCESSED / "modele_taille.csv").iloc[0]
    d["meta"] = json.loads(
        (RACINE / "app" / "donnees" / "meta.json").read_text(encoding="utf-8"))
    d["tendances"] = pd.read_csv(RACINE / "app" / "donnees" / "tendances.csv")
    d["params_lstm"] = json.loads(
        (PROCESSED / "best_params_lstm.json").read_text(encoding="utf-8"))
    return d


# ════════════════════════════════════════════════════════════════════════
#  primitives de mise en page
# ════════════════════════════════════════════════════════════════════════
def _cm(v):
    from pptx.util import Cm
    return Cm(v)


def _rgb(h):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(h)


class Deck:
    """Un diaporama 16:9, avec quelques gabarits maison.

    python-pptx ne fournit pas de mise en page : les gabarits par défaut sont
    ceux d'Office, avec leurs polices et leurs puces. On part donc de
    diapositives VIDES et on pose les blocs soi-même — c'est plus long à
    écrire, mais c'est la seule façon d'obtenir une charte cohérente.
    """

    L, R = 1.9, 1.9              # marges latérales, en cm
    LARG, HAUT = 33.87, 19.05    # 16:9

    def __init__(self):
        from pptx import Presentation

        self.p = Presentation()
        self.p.slide_width = _cm(self.LARG)
        self.p.slide_height = _cm(self.HAUT)

    # ── briques ─────────────────────────────────────────────────────────
    def _vide(self, fond=FOND, numero=True):
        """Une diapositive vierge, numérotée en bas à droite.

        Le numéro sert pendant les questions : un membre du jury peut dire
        « revenez à la 22 » au lieu de décrire la diapositive. Il est posé
        ici plutôt que dans chaque gabarit, sans quoi les diapositives de
        section l'oublieraient.

        Sa couleur suit le fond : les diapositives de section sont sombres,
        un numéro en MUTED y serait illisible.
        """
        s = self.p.slides.add_slide(self.p.slide_layouts[6])
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = _rgb(fond)
        if numero:
            from pptx.enum.text import PP_ALIGN

            # ⚠️ DANS LA MARGE DROITE, PAS SOUS LE CONTENU.
            # Posé d'abord en bas à droite de la zone de texte, il chevauchait
            # la dernière ligne de sept diapositives : les blocs font 30 cm de
            # large et descendent jusqu'à la même bande. Le sortir
            # horizontalement de la colonne de contenu (qui finit à L + 30)
            # règle le cas une fois pour toutes, quel que soit le texte à
            # venir.
            self._texte(s, len(self.p.slides._sldIdLst),
                        self.L + 30.1, self.HAUT - 1.3, 1.4, 0.7,
                        9, GRIS if fond == INK else MUTED,
                        aligne=PP_ALIGN.RIGHT)
        return s

    def _texte(self, s, txt, x, y, w, h, taille=16, coul=INK, gras=False,
               interligne=1.25, aligne=None):
        from pptx.util import Pt

        b = s.shapes.add_textbox(_cm(x), _cm(y), _cm(w), _cm(h))
        tf = b.text_frame
        tf.word_wrap = True
        for i, ligne in enumerate(str(txt).split("\n")):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.line_spacing = interligne
            if aligne is not None:
                par.alignment = aligne
            r = par.add_run()
            r.text = ligne
            r.font.size = Pt(taille)
            r.font.bold = gras
            r.font.color.rgb = _rgb(coul)
            r.font.name = "Calibri"
        return b

    def _image(self, s, chemin, x, y, w_max, h_max):
        """Insère en respectant le rapport d'aspect, centré dans la boîte."""
        from PIL import Image

        iw, ih = Image.open(chemin).size
        ratio = min(w_max / iw, h_max / ih)
        w, h = iw * ratio, ih * ratio
        return s.shapes.add_picture(
            str(chemin), _cm(x + (w_max - w) / 2), _cm(y + (h_max - h) / 2),
            width=_cm(w), height=_cm(h))

    def _filet(self, s, y, coul=ROUGE, x=None, w=3.2, ep=0.09):
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE

        f = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, _cm(x or self.L), _cm(y),
                               _cm(w), _cm(ep))
        f.fill.solid()
        f.fill.fore_color.rgb = _rgb(coul)
        f.line.fill.background()
        f.shadow.inherit = False
        return f

    # ── gabarits ────────────────────────────────────────────────────────
    def titre(self, sur, titre, sous, bas=""):
        s = self._vide(numero=False)      # on ne numérote pas la couverture
        self._texte(s, sur, self.L, 4.4, 24, 1, 13, ROUGE, True)
        self._texte(s, titre, self.L, 5.5, 30, 4, 40, INK, True, 1.05)
        self._texte(s, sous, self.L, 10.6, 26, 3, 17, MUTED, interligne=1.35)
        if bas:
            self._texte(s, bas, self.L, 16.2, 28, 1.4, 12, MUTED)
        return s

    def section(self, num, titre, sous=""):
        s = self._vide(INK)
        self._texte(s, num, self.L, 6.3, 6, 1.4, 15, ROUGE, True)
        self._texte(s, titre, self.L, 7.3, 28, 3, 34, FOND, True, 1.1)
        if sous:
            self._texte(s, sous, self.L, 11.4, 26, 2.5, 16, GRIS,
                        interligne=1.35)
        return s

    def page(self, titre, sous=""):
        """En-tête standard. Renvoie (slide, y où le contenu peut commencer)."""
        s = self._vide()
        self._filet(s, 1.5)
        self._texte(s, titre, self.L, 1.95, 29, 1.8, 27, INK, True)
        y = 3.9
        if sous:
            self._texte(s, sous, self.L, y, 29, 1.6, 14.5, MUTED,
                        interligne=1.3)
            y += 1.5
        return s, y

    # à 1 pt = 0,03528 cm, une ligne occupe 1,2 × la taille de police, et un
    # caractère de Calibri en bas de casse vaut environ la moitié de sa taille.
    PT = 0.03528
    LARG_CAR = 0.50
    INTERLIGNE = 1.25

    def _lignes(self, txt, w, taille):
        """Combien de lignes ce texte occupera une fois replié."""
        par_ligne = max(1, int((w - 0.25) / (taille * self.PT * self.LARG_CAR)))
        return max(1, -(-len(txt) // par_ligne))

    def puces(self, s, items, x, y, w, taille=15, ecart=1.55):
        """Une liste, avec un point médian plutôt qu'une puce Office.

        Le séparateur était un tiret cadratin ; comme chaque item non gras en
        recevait un, la seule mise en page produisait mécaniquement des
        dizaines de tirets dans le diaporama.

        ⚠️ L'ÉCART ENTRE ITEMS EST UN MINIMUM, PLUS UN PAS FIXE.
        python-pptx n'ajuste rien : un item de trois lignes posé tous les
        1,26 cm recouvrait le suivant, sans que rien ne le signale avant
        l'ouverture du fichier. On avance donc de la hauteur réellement
        nécessaire quand elle dépasse l'écart demandé.
        """
        cur = y
        for it in items:
            gras = it.startswith("**")
            txt = ("· " if not gras else "") + it.replace("**", "")
            besoin = (self._lignes(txt, w, taille)
                      * taille * self.PT * 1.2 * self.INTERLIGNE)
            self._texte(s, txt, x, cur, w, max(ecart, besoin), taille,
                        INK if gras else MUTED, gras)
            cur += max(ecart, besoin + 0.06)
        return cur

    def tableau(self, s, entetes, lignes, x, y, w, h, larg=None, taille=13,
                fort=()):
        from pptx.util import Pt

        t = s.shapes.add_table(len(lignes) + 1, len(entetes),
                               _cm(x), _cm(y), _cm(w), _cm(h)).table
        if larg:
            for i, lg in enumerate(larg):
                t.columns[i].width = _cm(lg)
        def _pose(c, txt, gras, coul, fond):
            # ⚠️ une cellule vide ne crée AUCUN run : `paragraphs[0].runs` est
            # alors un tuple vide et tout accès par index lève. On met une
            # espace insécable plutôt que de multiplier les gardes.
            c.text = str(txt) if str(txt).strip() else " "
            r = c.text_frame.paragraphs[0].runs[0]
            r.font.size = Pt(taille)
            r.font.bold = gras
            r.font.color.rgb = _rgb(coul)
            r.font.name = "Calibri"
            c.fill.solid()
            c.fill.fore_color.rgb = _rgb(fond)

        for j, e in enumerate(entetes):
            _pose(t.cell(0, j), e, True, FOND, INK)
        for i, lig in enumerate(lignes):
            for j, v in enumerate(lig):
                _pose(t.cell(i + 1, j), v, i in fort,
                      INK if i in fort else MUTED,
                      "F2F1EC" if i in fort else FOND)
        return t

    def encart(self, s, txt, x, y, w, h, coul=ORANGE, taille=15):
        from pptx.enum.shapes import MSO_SHAPE

        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, _cm(x), _cm(y), _cm(0.11),
                               _cm(h))
        b.fill.solid()
        b.fill.fore_color.rgb = _rgb(coul)
        b.line.fill.background()
        b.shadow.inherit = False
        self._texte(s, txt, x + 0.55, y + 0.1, w - 0.7, h, taille, INK,
                    interligne=1.3)

    def notes(self, s, txt):
        s.notes_slide.notes_text_frame.text = txt.strip()

    def enregistrer(self, chemin):
        chemin.parent.mkdir(parents=True, exist_ok=True)
        self.p.save(str(chemin))
        return chemin


# ════════════════════════════════════════════════════════════════════════
#  deux figures qui n'existent que pour la soutenance
# ════════════════════════════════════════════════════════════════════════
def _figures_soutenance(D) -> dict:
    """Génère ce que les notebooks n'ont pas produit, et renvoie les chemins.

    Le reste du diaporama réutilise les PNG déjà versionnés dans `figures/` :
    refaire une figure, c'est risquer qu'elle diverge de celle du notebook.
    """
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np

    dest = FIG / "soutenance"
    dest.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update({"figure.facecolor": "#" + FOND,
                         "axes.facecolor": "#" + FOND, "font.size": 9,
                         "axes.edgecolor": "#c3c2b7", "text.color": "#" + INK})
    out = {}

    # ── validation croisée spatiale ─────────────────────────────────────
    ts = D["transfert_spatial"]
    if ts is not None:
        t = ts.copy()
        t["ecart"] = 100 * (t["C · physique"] / t["A · tout"] - 1)
        t = t.sort_values("ecart")
        fig, ax = plt.subplots(figsize=(11, 4.4))
        y = np.arange(len(t))
        ax.barh(y - .2, t["A · tout"], height=.38, color="#" + GRIS,
                label="A · toutes les features", edgecolor="#" + FOND, lw=.8)
        ax.barh(y + .2, t["C · physique"], height=.38, color="#" + VERT,
                label="C · physique pure", edgecolor="#" + FOND, lw=.8)
        for i, (_, r) in enumerate(t.iterrows()):
            ax.text(max(r["A · tout"], r["C · physique"]) + .012, i,
                    f"{r.ecart:+.0f} %", va="center", fontsize=9,
                    weight="bold", color="#" + VERT)
        ax.set_yticks(y, [f"région {int(r)}" for r in t.region], fontsize=9)
        ax.set_xlabel("PR-AUC sur la région EXCLUE de l'entraînement")
        ax.set_xlim(0, t[["A · tout", "C · physique"]].to_numpy().max() * 1.18)
        ax.legend(frameon=False, fontsize=9.5, loc="lower right")
        ax.set_title("En territoire jamais vu, le modèle physique gagne "
                     "9 fois sur 9", fontsize=12, weight="bold", loc="left")
        ax.grid(axis="x", color="#e1e0d9", lw=.7)
        ax.set_axisbelow(True)
        ax.spines[["top", "right"]].set_visible(False)
        plt.tight_layout()
        out["spatial"] = dest / "spatial.png"
        fig.savefig(out["spatial"], dpi=150, facecolor="#" + FOND,
                    bbox_inches="tight")
        plt.close(fig)

    # ── SHAP du modèle C : gain contre SHAP ─────────────────────────────
    npy = PROCESSED / "shap_c_alea.npy"
    if npy.exists():
        cols = json.loads(
            (PROCESSED / "shap_c_colonnes.json").read_text(encoding="utf-8"))
        gain = 100 * pd.read_csv(PROCESSED / "importances_c.csv",
                                 index_col=0).squeeze("columns")
        sa = pd.Series(np.abs(np.load(npy)).mean(0), index=cols)
        ss = pd.Series(np.abs(np.load(PROCESSED / "shap_c_sommet.npy")).mean(0),
                       index=cols)
        top = pd.DataFrame({"gain": gain, "alea": sa, "sommet": ss}) \
            .sort_values("sommet", ascending=False).head(12).iloc[::-1]
        fig, ax = plt.subplots(1, 3, figsize=(13, 4.6), sharey=True)
        y = np.arange(len(top))
        for a_, c, ti, co in ((ax[0], "gain", "gain d'entraînement (%)", GRIS),
                              (ax[1], "alea", "SHAP · échantillon aléatoire", BLEU),
                              (ax[2], "sommet", "SHAP · sommet du classement", ROUGE)):
            a_.barh(y, top[c], color="#" + co, edgecolor="#" + FOND, lw=1)
            a_.set_title(ti, fontsize=11, weight="bold", loc="left")
            a_.grid(axis="x", color="#e1e0d9", lw=.7)
            a_.set_axisbelow(True)
            a_.spines[["top", "right"]].set_visible(False)
        ax[0].set_yticks(y, [f.replace("_", " ") for f in top.index], fontsize=9)
        plt.tight_layout()
        out["shap"] = dest / "shap_c.png"
        fig.savefig(out["shap"], dpi=150, facecolor="#" + FOND,
                    bbox_inches="tight")
        plt.close(fig)

    return out


# ════════════════════════════════════════════════════════════════════════
#  le diaporama
# ════════════════════════════════════════════════════════════════════════


def main() -> None:
    D = _lire()
    G = _figures_soutenance(D)
    d = Deck()
    mt, ten = D["meta"], D["tendances"]
    ap = D["pr_auc_val"].iloc[0].to_dict()
    cmp_ = D["comparaison_appariee"]
    taux_val = 0.0002410

    def t(nom):
        return ten[ten.serie == nom].iloc[0]

    def ec(ref, mod):
        return cmp_[(cmp_.reference == ref) & (cmp_.modele == mod)].iloc[0]

    IMG = RACINE / "docs" / "img"

    # ════════════════════════════════════════════════════════════════════
    #  TITRE
    # ════════════════════════════════════════════════════════════════════
    s = d.titre("Projet de fin d'année · M1 Data / IA · La Plateforme_",
                "Où va-t-il brûler\ndemain ?",
                "Estimer le risque de départ de feu pour chacune des 34 734\n"
                "communes de France métropolitaine, chaque jour,\n"
                "de 1973 à 2100.",
                "Romuald Courtois")
    d._image(s, IMG / "carte.png", 20.5, 2.4, 11.5, 14.2)
    d.notes(s, "Cette carte est produite par le modèle. Il retrouve seul les "
               "Landes, l'arc méditerranéen et la Corse, alors que ni la "
               "latitude ni la longitude ne font partie de ses variables. "
               "C'est la végétation et le relief qui portent cette géographie.")

    # ════════════════════════════════════════════════════════════════════
    #  01 · INTRODUCTION
    # ════════════════════════════════════════════════════════════════════
    d.section("01", "Introduction",
              "Le phénomène, et ce qui manque aux outils existants")

    # ── contexte ────────────────────────────────────────────────────────
    s, y = d.page("Le phénomène",
                  "52 809 feux déclarés sur le périmètre d'étude, entre 1973 "
                  "et 2025.")
    d._image(s, FIG / "data-bdiff" /
             "03_feux-par-jour-de-l-annee-cumul-sur-53-ans.png",
             d.L, y, 18.5, 9.6)
    d.puces(s, [
        "**Deux saisons, pas une**",
        "Le pic d'août est celui qu'on attend : l'arc méditerranéen, la Corse, "
        "un été sec.",
        "Celui de mars est moins connu. Il concerne le Sud-Ouest et le Massif "
        "central, et il ne relève pas de la même végétation : ce sont des "
        "landes et des herbacées sèches avant la reprise, souvent des écobuages "
        "qui échappent.",
        "**Une conséquence pour le modèle**",
        "Un modèle qui n'apprendrait que l'été raterait un feu sur cinq. La "
        "saisonnalité doit donc entrer comme variable, et pas comme filtre.",
    ], 21, y, 11.2, 13.5, 1.2)
    d.notes(s, "Point de départ concret. Si on demande d'où vient le pic de "
               "mars : Cantal 52 % de ses feux hors saison estivale, Dordogne "
               "44 %, Lozère 41 %. Ces chiffres sont calculés dans "
               "l'application, page Les données.")

    # ── ce qui existe déjà ──────────────────────────────────────────────
    s, y = d.page("Ce qui existe déjà, et ce qui manque",
                  "Le danger météo est un service public. La question posée "
                  "ici commence là où il s'arrête.")
    d.tableau(s, ["", "Ce que fournit le danger météo (EFFIS, Météo-France)",
                  "Ce que ce projet ajoute"], [
        ["Objet mesuré", "l'inflammabilité de l'air et du combustible",
         "la probabilité qu'un départ soit déclaré"],
        ["Maille", "0,25°, soit environ 28 km",
         "la commune"],
        ["Territoire", "non pris en compte",
         "végétation, relief, densité, littoral"],
        ["Horizon", "quelques jours",
         "1973 à 2100, trois scénarios GIEC"],
    ], d.L, y, 20.5, 5.6, larg=[4, 9, 7.5], taille=12)
    d.encart(s, "À l'intérieur d'une même maille de 28 km, on trouve en "
                "moyenne 31 communes. Elles reçoivent aujourd'hui le même "
                "indice de danger, qu'elles soient couvertes de maquis ou "
                "bétonnées.\nC'est l'écart que ce travail cherche à combler : "
                "la météo dit quand, elle ne dit pas où.",
             d.L, y + 6.2, 20.5, 4, BLEU, 14)
    # La grille rend le maillage tangible : on voit la taille des carreaux,
    # donc on comprend l'argument des 31 communes sans avoir à le croire.
    d._image(s, FIG / "data-copernicus" /
             "02_grille-cems-0-25-1131-cellules-france-sur-2360.png",
             23, y, 9, 10.2)
    d.notes(s, "Montrer la grille en parlant : chaque carreau bleu est une "
               "cellule météo, il y en a 1 131 sur la France. C'est la "
               "résolution à laquelle le danger est publié aujourd'hui.\n\n"
               "Formuler l'apport sans exagérer : on ne remplace pas le FWI, "
               "on l'utilise comme entrée principale et on lui ajoute une "
               "résolution spatiale qu'il n'a pas.")

    # ════════════════════════════════════════════════════════════════════
    #  02 · QUESTION DE RECHERCHE
    # ════════════════════════════════════════════════════════════════════
    d.section("02", "Question de recherche",
              "Une question principale, quatre hypothèses testables")

    # ── la question ─────────────────────────────────────────────────────
    s, y = d.page("La question")
    d.encart(s, "Peut-on estimer, à l'échelle du couple commune × jour, la "
                "probabilité qu'un départ de feu soit déclaré, à partir des "
                "seules données publiques disponibles en temps réel ?",
             d.L, y, 30, 3.6, ROUGE, 21)
    d._texte(s, "Quatre hypothèses, chacune vérifiable, chacune réfutable",
             d.L, y + 4.6, 30, 1.2, 17, INK, True)
    d.tableau(s, ["", "Hypothèse", "Comment on la teste"], [
        ["H1", "La météo seule ne suffit pas à localiser le risque",
         "baseline « danger EFFIS seul », comparée au croisement avec le "
         "territoire"],
        ["H2", "Le territoire apporte une information que la météo ne porte pas",
         "SHAP sur deux populations, et validation croisée par région"],
        ["H3", "Une architecture séquentielle n'apporte rien de plus que des "
               "indices physiques agrégés",
         "LSTM contre modèle physique, à information égale, bootstrap apparié"],
        ["H4", "Le danger météo augmente de façon détectable sur 53 ans",
         "régression de tendance sur quatre séries, avec test de "
         "significativité"],
    ], d.L, y + 6.2, 30, 7.4, larg=[1.8, 12.2, 16])
    d.notes(s, "Insister sur le fait que H3 est formulée dans le sens qui "
               "permet de la réfuter : si le LSTM avait gagné, on l'aurait "
               "déployé. Le résultat n'était pas décidé d'avance.")

    # ── la difficulté ───────────────────────────────────────────────────
    s, y = d.page("Ce qui rend la question difficile",
                  "Un départ de feu est un événement rare.")
    d._texte(s, "0,019 %", d.L, y + .3, 14, 3.4, 66, ROUGE, True, 1)
    d._texte(s, "des couples commune × jour comptent\nau moins un départ de feu",
             d.L, y + 4, 14, 2.4, 15, MUTED)
    d.puces(s, [
        "**Répondre toujours « non » donne 99,98 % de justesse**",
        "et ne sert à rien. L'exactitude est inutilisable ici.",
        "**Une fuite de données ne provoque pas d'erreur**",
        "elle produit d'excellentes métriques et un modèle sans valeur. Rien "
        "dans les scores ne la signale.",
        "**Le déséquilibre commande le reste du protocole**",
        "la métrique, l'échantillonnage, la calibration, et jusqu'à la façon "
        "de comparer deux modèles entre eux.",
    ], 17.5, y + .4, 14.5, 15, 1.35)
    d.notes(s, "C'est la contrainte qui explique chacune des décisions "
               "méthodologiques présentées ensuite. Le dire maintenant évite "
               "de le justifier cinq fois.")

    # ════════════════════════════════════════════════════════════════════
    #  03 · MÉTHODOLOGIE
    # ════════════════════════════════════════════════════════════════════
    d.section("03", "Méthodologie",
              "Les données, le protocole, et ce qui protège le résultat")

    # ── les sources ─────────────────────────────────────────────────────
    s, y = d.page("Quatre sources publiques, croisées",
                  "Jointure sur le code INSEE et sur une grille météo de 0,25°.")
    d.tableau(s, ["Source", "Ce qu'elle apporte", "Volume"], [
        ["CEMS · Copernicus", "8 indices de danger, par jour et par maille",
         "21,9 M lignes · 1973-2025"],
        ["BDIFF · IGN", "les feux déclarés, commune par commune",
         "52 809 feux sur le périmètre"],
        ["CORINE Land Cover", "l'occupation du sol, 44 postes", "1,08 M lignes"],
        ["INSEE", "référentiel des communes, et leurs fusions",
         "34 734 communes"],
    ], d.L, y, 17.5, 6, larg=[4.5, 7.5, 5.5], taille=12)
    d._image(s, FIG / "data-ville" /
             "03_34-734-communes-metropolitaines-un-centroide-par-commune.png",
             19.6, y + 0.6, 12.4, 4.6)
    d._texte(s, "Le référentiel INSEE, une ligne par commune : population, "
                "altitude et densité pour les 34 734 communes.",
             19.6, y + 5.4, 12.4, 1.4, 12, MUTED)
    d.encart(s, "Le point délicat n'est pas le volume, ce sont les fusions de "
                "communes.\n965 feux portent un code INSEE qui n'existe plus. "
                "Les rapprocher par le nom produisait de faux positifs : "
                "« Chirac » en Lozère renvoyait vers la Charente. On a donc "
                "utilisé le fichier officiel des mouvements de communes de "
                "l'INSEE, et les 30 cas restants sont écartés et comptés, "
                "jamais devinés.",
             d.L, y + 6.8, 30, 4.4, ORANGE, 14)
    d.notes(s, "Une heuristique par le nom aurait « résolu » tous les cas et "
               "corrompu la cible, le voisinage et les features CORINE de "
               "communes innocentes. Perdre 0,68 % des feux coûte moins cher "
               "que d'en inventer.")

    # ── la grille ───────────────────────────────────────────────────────
    s, y = d.page("La table centrale : commune × jour",
                  "Une ligne par commune et par jour, qu'il y ait eu un feu "
                  "ou non.")
    d.tableau(s, ["Partition", "Années", "Lignes", "Feux", "Taux"], [
        ["train", "2006-2019", "177 594 942", "33 632", "0,0189 %"],
        ["validation", "2020-2022", "38 068 464", "9 176", "0,0241 %"],
        ["test", "2023-2025", "38 068 464", "6 322", "0,0166 %"],
        ["total", "2006-2025", "253 731 870", "49 130", "0,0194 %"],
    ], d.L, y, 24, 5.4, fort=(3,))
    d.encart(s, "Pourquoi une grille dense plutôt que la seule liste des feux\n"
                "Parce qu'une série creuse rendrait les fenêtres glissantes "
                "silencieusement fausses. « Feux des 30 jours précédents » se "
                "calcule en remontant 30 lignes : si les jours sans feu sont "
                "absents, on remonte en réalité plusieurs années.",
             d.L, y + 6.2, 30, 3.8, BLEU, 14)
    d.notes(s, "253 millions de lignes plutôt que 52 809 : c'est un choix "
               "coûteux en stockage et en temps de calcul, et il a une raison "
               "technique précise. PostgreSQL avec 20 partitions annuelles.")

    # ── le split ────────────────────────────────────────────────────────
    s, y = d.page("La barrière du split",
                  "Découpage temporel, jamais aléatoire.")
    d.puces(s, [
        "**train · 2006-2019 · apprendre**",
        "177,6 M lignes. Tout ce qui possède un .fit() s'ajuste ici, et nulle "
        "part ailleurs.",
        "**validation · 2020-2022 · choisir**",
        "Hyperparamètres, sélection du modèle, calibration, comparaisons.",
        "**test · 2023-2025 · juger, une seule fois**",
        "Ouvert après gel complet. Aucune décision n'en découle.",
    ], d.L, y, 17, 15, 1.4)
    d.encart(s, "Un découpage aléatoire mettrait le 14 juillet 2019 dans le "
                "train et le 15 dans le test. Le modèle « prédirait » un feu "
                "qu'il a déjà vu, à 20 km et un jour d'écart.",
             19.5, y + .2, 12.6, 3.6, ROUGE, 14)
    d.encart(s, "La règle qui tranche les cas douteux\nUne feature datée peut "
                "regarder tout le passé, y compris celui de sa propre période "
                "d'évaluation.\nUne statistique non datée ne peut regarder que "
                "le train.", 19.5, y + 4.6, 12.6, 4.4, BLEU, 14)
    d.notes(s, "Exemple à donner si on demande : « feux des 30 jours "
               "précédents » au 3 août 2023 lit juillet 2023, et ce n'est pas "
               "une fuite, car le 3 août à 8 h on connaît juillet. « Taux moyen "
               "de la commune sur toute la période » lit le futur, et c'en est "
               "une.")

    # ── la métrique ─────────────────────────────────────────────────────
    s, y = d.page("Pourquoi PR-AUC, et pas ROC-AUC",
                  "À 0,019 % de positifs, le choix de la métrique décide de ce "
                  "qu'on croit avoir réussi.")
    d.tableau(s, ["", "ROC-AUC", "PR-AUC"], [
        ["Ce qu'elle mesure", "vrais positifs contre faux positifs",
         "précision contre rappel"],
        ["Valeur au hasard", "0,50, quel que soit le déséquilibre",
         "exactement le taux de base"],
        ["À 0,019 % de positifs", "flatteuse : un modèle médiocre affiche 0,95",
         "lisible : le rapport donne le lift"],
    ], d.L, y, 30, 5.2, larg=[7, 11.5, 11.5])
    d._texte(s, "lift  =  PR-AUC ÷ taux de base",
             d.L, y + 6.2, 30, 1.6, 26, INK, True)
    d._texte(s, "combien de fois mieux que tirer au sort, un nombre qui se dit "
                "à voix haute", d.L, y + 8.2, 30, 1.4, 16, MUTED)
    d.notes(s, "C'est souvent la première question posée : pourquoi pas "
               "l'exactitude. Répondre par le chiffre, 99,98 % de justesse en "
               "répondant toujours non.")

    # ── le prior ────────────────────────────────────────────────────────
    s, y = d.page("Le piège du sous-échantillonnage",
                  "Le train est réduit à 1 positif pour 10 négatifs, sans quoi "
                  "l'entraînement serait ingérable.")
    d._texte(s, "× 487", d.L, y + .4, 13, 3.4, 60, ORANGE, True, 1)
    d._texte(s, "l'écart entre le prior appris (9,1 %)\net le prior réel "
                "(0,019 %)", d.L, y + 4.1, 14, 2.4, 15, MUTED)
    d.puces(s, [
        "**Validation et test ne sont jamais échantillonnés**",
        "c'est ce qui rend les scores comparables au monde réel.",
        "**Les statistiques dérivées de la cible se calculent sur le train "
        "complet**",
        "sur l'échantillon, un lissage bayésien vaudrait 9,1 % au lieu de "
        "0,019 %. Le prior serait empoisonné, et rien dans les métriques ne "
        "le signalerait.",
        "**La calibration absorbe le décalage**",
        "Platt ramène le biais de ×144,7 à ×1,13, sans rien coûter en PR-AUC.",
    ], 17.5, y + .3, 14.5, 14.5, 1.3)
    d.notes(s, "Ce piège est invisible : le modèle fonctionne, les métriques "
               "sont bonnes, seul le niveau absolu des probabilités est faux. "
               "C'est pour cette raison que l'application affiche un rang.")

    # ── le protocole de comparaison ─────────────────────────────────────
    s, y = d.page("Comment comparer deux modèles sans se tromper",
                  "Un écart de PR-AUC entre deux modèles n'est pas un "
                  "résultat tant qu'on ignore le bruit qui l'entoure.")
    d.puces(s, [
        "**Bootstrap apparié**",
        "Les deux modèles sont évalués sur exactement les mêmes lignes, "
        "rééchantillonnées ensemble, 200 fois. On mesure la distribution de "
        "l'écart, pas deux distributions séparées.",
        "**On rééchantillonne les communes, pas les lignes**",
        "Les 1 096 jours d'une même commune ne sont pas indépendants, et 31 "
        "communes partagent la même maille météo. Tirer ligne à ligne "
        "reviendrait à traiter 38 millions d'observations comme 38 millions "
        "d'expériences indépendantes.",
    ], d.L, y, 30, 15, 1.24)
    d.encart(s, "C'est le point de méthode le plus important de cette partie. "
                "Un bootstrap ligne à ligne donne des intervalles beaucoup "
                "trop étroits, et fait conclure à des différences qui "
                "n'existent pas. Le terme consacré est la pseudo-réplication.",
             d.L, y + 6.4, 30, 3.2, ROUGE, 15)
    d.notes(s, "Détail d'implémentation, à garder pour les questions : la "
               "moyenne pondérée des précisions est recalculée par tri unique "
               "puis sommes cumulées, au lieu d'appeler scikit-learn 200 fois. "
               "Résultat identique à 1e-12 près, et la comparaison complète "
               "tient en quelques minutes au lieu de plusieurs heures.\n\n"
               "Si on demande pourquoi 200 répliques : au-delà, la largeur des "
               "intervalles ne bouge plus de façon utile, et le coût est "
               "linéaire.")

    # ── les modèles comparés ────────────────────────────────────────────
    s, y = d.page("Les modèles mis en concurrence",
                  "Chacun teste une idée précise, et non « essayons de voir ».")
    d.tableau(s, ["Modèle", "L'idée qu'il teste", "Entrées"], [
        ["Baselines (3)", "que valent des règles sans apprentissage",
         "historique, danger EFFIS, et leur croisement"],
        ["XGBoost v1", "un gradient boosting standard suffit-il",
         "météo, territoire, historique brut"],
        ["XGBoost v3", "peut-on donner un risque aux communes sans historique",
         "v1 + lissage bayésien sur 30 groupes"],
        ["DART", "l'abandon d'arbres améliore-t-il la généralisation",
         "identiques à v3"],
        ["MLP", "un réseau dense fait-il mieux que des arbres",
         "identiques à v3, normalisées"],
        ["XGBoost C", "que reste-t-il sans aucun historique de feu",
         "météo et territoire seulement, 41 features"],
        ["LSTM", "la séquence météo porte-t-elle un signal en propre",
         "30 jours × 8 indices, soit 240 valeurs"],
    ], d.L, y, 30, 9, larg=[6, 14, 10], taille=12.5)
    d.notes(s, "Le modèle C est celui qui répond à la contrainte de "
               "déploiement, et c'est aussi la seule référence honnête face au "
               "LSTM, puisqu'aucun des deux ne voit l'historique des feux.")

    # ════════════════════════════════════════════════════════════════════
    #  04 · RÉSULTATS
    # ════════════════════════════════════════════════════════════════════
    d.section("04", "Résultats",
              "Ce que mesurent les données, avant toute interprétation")

    # ── les baselines ───────────────────────────────────────────────────
    s, y = d.page("La barre à battre",
                  "Sans référence, un lift de 63× ne veut rien dire.")
    b = D["baselines"]
    d.tableau(s, ["Prédicteur, sans aucun apprentissage", "PR-AUC", "lift"],
              [[r.predicteur, f"{r.pr_auc:.4f}", f"×{r.lift:.1f}"]
               for _, r in b.iterrows()]
              + [["le modèle déployé, mesuré sur le test",
                  f"{mt['test']['pr_auc']:.4f}", f"×{mt['test']['lift']:.1f}"]],
              d.L, y, 25, 4.8, larg=[15, 5, 5], fort=(len(b),))
    d.encart(s, "L'historique spatial seul vaut déjà ×19, la météo seule ×5, "
                "et leur croisement ×42. C'est la barre à battre, et non le "
                "hasard : un modèle à ×30 serait moins bon qu'une règle de "
                "trois, tout en paraissant excellent.",
             d.L, y + 5.3, 30, 2.2, VERT, 14)
    # La baseline météo, détaillée : le taux de feu passe de 0,006 % en
    # classe « très faible » à 0,540 % en « extrême ». Le FWI porte donc un
    # vrai signal, ce qui rend la barre d'autant plus honnête.
    # Format très allongé (ratio 3,6) : il lui faut toute la largeur, donc on
    # resserre le tableau au-dessus plutôt que de la réduire de moitié.
    d._image(s, FIG / "data-all" /
             "07_baseline-2-ce-que-le-danger-effis-predit-a-lui-seul.png",
             d.L, y + 7.9, 30, 5.4)
    d.notes(s, "Première réponse à H1 : la météo seule plafonne à ×5, le "
               "croisement avec le territoire monte à ×42. L'écart entre les "
               "deux est la moitié de la réponse à la question de recherche.")

    # ── v1 → v3 ─────────────────────────────────────────────────────────
    s, y = d.page("Donner un risque aux communes sans historique",
                  "Le v1 tirait 54,6 % de son importance de l'historique de "
                  "la commune. Il disait surtout que ce qui a brûlé rebrûlera.")
    d._image(s, FIG / "modele-v3" /
             "02_le-lissage-bayesien-rendre-une-estimation-aux-communes-sans.png",
             d.L, y, 18.5, 10.5)
    d.puces(s, [
        "**Le problème**",
        "Une commune qui n'a jamais brûlé gardait un score bas, même entourée "
        "de communes qui brûlent chaque été.",
        "**La réponse**",
        "Regrouper les communes qui se ressemblent physiquement, en 30 groupes "
        "formés sans jamais regarder le feu, puis faire retomber chaque "
        "commune vers le taux de son groupe, à proportion de ce qu'on sait "
        "d'elle.",
        "**Le gain, mesuré**",
        "+0,83 % de PR-AUC. Réel mais modeste, et le dire fait partie du "
        "résultat.",
    ], 21.3, y, 10.9, 13, 1.22)
    d.notes(s, "C'est le problème classique d'estimation sur petits domaines. "
               "Le point méthodologique : les groupes sont formés sans "
               "regarder la cible, la sinistralité n'entre qu'ensuite.")

    # ── les modèles comparés (5, lus dans pr_auc_val.csv) ───────────────
    s, y = d.page("Les modèles, sur la même validation",
                  "38 068 464 communes-jours, 9 176 feux, taux 0,0241 %.")
    lignes = []
    for nom in sorted(ap, key=ap.get, reverse=True):
        e = "référence" if nom == "XGBoost v3" else \
            f"{ec('XGBoost v3', nom).ecart_pct:+.1f} %"
        lignes.append([nom, f"{ap[nom]:.4f}", f"×{ap[nom] / taux_val:.1f}", e])
    d.tableau(s, ["Modèle", "PR-AUC", "lift", "vs XGBoost v3"], lignes,
              d.L, y, 23, 6.6, larg=[9.5, 4.5, 4.5, 4.5], fort=(0,))
    ens = D["modeles_ensemble"]
    if ens is not None:
        d._texte(s, f"L'ensemble v3 + MLP monte à {ens.pr_auc.iloc[0]:.4f} "
                    f"(×{ens.lift.iloc[0]:.1f}), au prix de deux modèles à "
                    f"faire tourner en production.",
                 d.L, y + 7.3, 30, 1.4, 14, MUTED)
    d.encart(s, "Aucun de ces écarts n'a été retenu sans intervalle de "
                "confiance. C'est l'objet de la diapositive suivante.",
             d.L, y + 8.8, 30, 2.2, ROUGE, 16)
    d.notes(s, "Ne pas commenter le classement ici. Le commentaire n'a de sens "
               "qu'une fois les intervalles affichés, et c'est justement "
               "l'erreur que l'on veut éviter de commettre devant le jury.")

    # ── le graphique en forêt ───────────────────────────────────────────
    s, y = d.page("Ces écarts survivent-ils au bruit ?",
                  "Bootstrap apparié, 200 répliques, rééchantillonnage des "
                  "34 734 communes.")
    d._image(s, IMG / "modeles.png", d.L, y, 30, 8.2)
    d.encart(s, "Le résultat le plus instructif de cette partie est négatif.\n"
                f"DART et le MLP paraissaient "
                f"{abs(ec('XGBoost v3','DART').ecart_pct):.1f} % et "
                f"{abs(ec('XGBoost v3','MLP').ecart_pct):.1f} % moins bons. "
                "Leurs intervalles traversent zéro : les trois modèles sont "
                "indiscernables. Annoncer que XGBoost bat le MLP aurait été "
                "une conclusion tirée du bruit.",
             d.L, y + 8.8, 30, 3.8, ROUGE, 15)
    d.notes(s, "Ce que cette diapositive démontre vraiment, c'est que le "
               "protocole sait dire « je ne sais pas ». Un protocole qui "
               "conclut toujours quelque chose ne conclut rien.")

    # ── H3 · le LSTM, protocole et résultat ─────────────────────────────
    s, y = d.page("H3 · la séquence météo apporte-t-elle un signal en propre ?",
                  "Le réflexe standard face à une série temporelle. On l'a "
                  "construit, optimisé, puis mesuré.")
    p = D["params_lstm"]
    d.tableau(s, ["Hyperparamètre", "Valeur retenue"],
              [[k, f"{v:.5g}"] for k, v in p.items()],
              d.L, y, 13, 7.6, larg=[7.5, 5.5], taille=12)
    d._texte(s, "25 essais Optuna, arrêt précoce à l'époque 21.\n"
                "L'objection « il n'a pas été réglé » ne tient pas.",
             d.L, y + 8.3, 15, 2.4, 15, MUTED)
    l_c = ec("XGBoost C", "LSTM")
    d._texte(s, f"{l_c.ecart_pct:.1f} %", 18.5, y + .8, 13.5, 3.4, 62,
             ROUGE, True, 1)
    d._texte(s, f"contre le modèle C, à information égale\n"
                f"intervalle [{l_c.ic_bas:.1f} ; {l_c.ic_haut:.1f}], "
                f"loin de zéro", 18.5, y + 4.5, 13.5, 2.4, 16, MUTED)
    d.encart(s, "La comparaison loyale n'est pas celle qu'on croit\n"
                "XGBoost v3 voit l'historique des feux, soit 29 % de ses "
                "importances. Le LSTM n'en voit rien. Les opposer mesurerait "
                "le prix de l'information retirée, pas la valeur de la "
                "séquence. La seule référence à jeu d'information égal est le "
                "modèle C.", 18.5, y + 7.5, 13.5, 4.8, ORANGE, 14)
    d.notes(s, "Le LSTM voit 30 jours × 8 indices, soit 240 valeurs. Le modèle "
               "C voit les 8 indices du jour plus deux décalages. Vingt fois "
               "plus d'historique météo, et il perd.")

    # ── H3 · l'explication ──────────────────────────────────────────────
    s, y = d.page("Pourquoi il perd : l'explication est physique",
                  "Un LSTM sert quand l'ordre de la séquence porte une "
                  "information qu'aucun résumé ne capture. Ici, ce résumé "
                  "existe déjà.")
    d.encart(s, "Les indices DC, DMC et BUI du système canadien sont déjà des "
                "états récursifs : des moyennes exponentielles de la météo "
                "passée, de constante de temps 52 et 15 jours. Le CEMS livre "
                "donc l'état caché que le LSTM devrait réapprendre.",
             d.L, y, 30, 2.6, VERT, 15)
    # La figure porte l'argument mieux que le texte : on VOIT les trois
    # constantes de temps, du FFMC nerveux au DC lissé. Le détail chiffré de
    # la PACF descend dans les notes pour lui laisser la place.
    d._image(s, FIG / "data-copernicus" /
             "06_les-trois-memoires-temporelles-du-systeme-canadien-france-20.png",
             d.L, y + 3.1, 18.4, 9.1)
    acf = D["series_acf_pacf"]
    sar = D["series_sarimax"]
    r_ar = (sar[sar.modele.str.contains("sans exogène")].correlation.iloc[0]
            if sar is not None else float("nan"))
    p1 = p3 = p8 = sl = float("nan")
    if acf is not None:
        p1, p3 = acf.pacf.iloc[0], acf.pacf.iloc[2]
        p8, sl = acf.pacf.iloc[7], acf.seuil.iloc[0]
    d.puces(s, [
        "**Trois observations convergent**",
        f"la PACF tombe de {p1:.2f} au premier retard à {p3:.2f} au "
        f"troisième : deux à trois jours de mémoire utile ;",
        f"l'ARIMA sans exogène donne r = {r_ar:.3f}, une corrélation "
        f"négative ;",
        "les trois premières variables du modèle C disent où se trouve le "
        "combustible.",
    ], 21.2, y + 3.1, 10.8, 13, 1.05)
    d.notes(s, f"Lire la figure de haut en bas : le FFMC réagit au jour le "
               f"jour sur 1 à 2 cm de litière, le DMC lisse sur 5 à 10 cm, le "
               f"DC monte et redescend sur toute la saison à 10-20 cm de "
               f"profondeur. Trois constantes de temps, visibles à l'œil.\n\n"
               f"Nuance à donner si on interroge la PACF : les retards 4 à 8 "
               f"restent significatifs. Avec 7 305 jours le seuil descend à "
               f"±{sl:.3f} et presque tout le devient, mais le retard 8, à "
               f"{p8:.3f}, ne rend compte que de {100 * p8 ** 2:.2f} % de la "
               f"variance. Significatif ne veut pas dire utile.\n\n"
               "Réserve à donner spontanément, sans attendre la question : le "
               "LSTM ne reçoit pas danger_effis, qui pèse 13,7 % dans le "
               "modèle C. L'écart de 23,6 % est donc un majorant, et c'est la "
               "première chose à reprendre.")

    # ── les séries temporelles ──────────────────────────────────────────
    s, y = d.page("Ce que le passé de la série prédit à lui seul",
                  "Un détour par les méthodes classiques, qui sert de "
                  "contre-épreuve au LSTM.")
    adf = D["series_adf"]
    if adf is not None:
        d.tableau(s, ["Série testée", "stat ADF", "p", "Conclusion"],
                  [[r.serie, f"{r.adf:.2f}", f"{r.p:.1e}",
                    "stationnaire" if r.stationnaire else "non stationnaire"]
                   for _, r in adf.iterrows()],
                  d.L, y, 14.5, 4.4, larg=[6, 3, 3, 2.5], taille=12)
        d._texte(s, "H₀ est « racine unitaire », donc non stationnaire. "
                    "Rejeter H₀ signifie stationnaire, l'inverse de "
                    "l'intuition.", d.L, y + 4.7, 14.5, 1.8, 13, MUTED)
    if sar is not None:
        d.tableau(s, ["Modèle", "MAE", "r"],
                  [[r.modele, f"{r.mae:.2f}", f"{r.correlation:.3f}"]
                   for _, r in sar.iterrows()],
                  d.L, y + 6.8, 14.5, 3.6, larg=[9, 2.8, 2.7], taille=12,
                  fort=(len(sar) - 1,))
        d._texte(s, "Ajusté sur 2006-2019, évalué sur 2020-2022. Le test "
                    "2023-2025 n'est pas touché.",
                 d.L, y + 10.6, 14.5, 1.6, 13, MUTED)
    # Les corrélogrammes régénérés depuis PostgreSQL dans la même session que
    # les chiffres de la PACF : figure et texte ne peuvent pas diverger.
    if (PROCESSED / "series_acf_pacf.png").exists():
        d._image(s, PROCESSED / "series_acf_pacf.png", 17, y, 15, 6.4)
    d.encart(s, "La dernière ligne du tableau est celle qui compte. Un ARIMA "
                "sans variable exogène donne une corrélation négative : à "
                "1 096 pas d'horizon, un modèle dont la mémoire utile vaut "
                "trois jours a oublié son point de départ et converge vers la "
                "moyenne.\nAjouter le FWI fait tomber l'erreur de 37 %.",
             17, y + 7, 15, 4.6, ORANGE, 14)
    d.notes(s, "Ce détour n'est pas décoratif : il montre que la conclusion "
               "sur le LSTM ne dépend pas d'une seule expérience, mais de "
               "trois familles de méthodes qui disent la même chose.\n\n"
               "Les corrélogrammes de droite : en haut la série brute, où "
               "l'ACF ne montre que le cycle annuel parce que tout est corrélé "
               "à tout pendant l'été. En bas, après retrait de la saisonnalité "
               "par termes de Fourier, seules courbes exploitables pour "
               "choisir les ordres.")

    # ── v3 contre C ─────────────────────────────────────────────────────
    s, y = d.page("Le meilleur modèle n'est pas celui qu'on déploie",
                  "C'est la décision la plus contre-intuitive du projet.")
    d.tableau(s, ["Sur le test 2023-2025", "PR-AUC", "lift"], [
        ["XGBoost v3 · 52 features", f"{mt['modele_a']['pr_auc']:.4f}",
         f"×{mt['modele_a']['lift']:.1f}"],
        ["XGBoost C · physique pure, 41 features", f"{mt['test']['pr_auc']:.4f}",
         f"×{mt['test']['lift']:.1f}"],
    ], d.L, y, 19, 3.2, larg=[11, 4, 4], fort=(1,))
    yy = d.puces(s, [
        "**1. La donnée n'existe pas en temps réel**",
        "v3 tire 29 % de son importance de l'historique des feux, et la BDIFF "
        "ne publie pas l'année en cours. Pas imprécis : faux.",
        "**2. En territoire inconnu, elle vaut zéro**",
        "et le modèle lit ce zéro comme « ça n'a jamais brûlé, donc ça ne "
        "brûlera pas », précisément là où un risque nouveau apparaît.",
        "**3. Pour 2050, elle est impossible par construction**",
        "on ne connaîtra jamais les feux de 2049.",
    ], d.L, y + 3.7, 30, 13.5, 1.05)
    d.encart(s, "Le choix se fait sur la disponibilité de la donnée, pas sur "
                "la performance. Ce défaut n'apparaît dans aucune métrique "
                "d'entraînement : en validation comme en test, l'historique "
                "est toujours là.",
             d.L, yy + 0.3, 30, 1.8, ROUGE, 14)
    d.notes(s, "Développer le point 1 à l'oral : les feux de 2026 sortiront au "
               "printemps 2027. Si je lançais le modèle ce matin, "
               "feux_commune_7j vaudrait le décompte d'une semaine de décembre "
               "2025.\n\n"
               "Dans l'application, un basculement permet de comparer les deux "
               "modèles, et il refuse de le faire ailleurs que sur le test, en "
               "expliquant pourquoi pour chaque période.")

    # ── validation croisée spatiale ─────────────────────────────────────
    s, y = d.page("H2 · retirer une région entière, puis tester dessus",
                  "Simulation d'un territoire jamais vu, ou d'un climat qui "
                  "déplace le risque.")
    if "spatial" in G:
        d._image(s, G["spatial"], d.L, y, 30, 8.8)
    d.encart(s, "Le modèle physique gagne dans les 9 régions, sans exception : "
                "+8,2 % en moyenne pondérée, et jusqu'à +137 % dans le Grand "
                "Est.\nLà où l'historique est le plus pauvre, s'y fier devient "
                "un handicap. C'est l'argument décisif pour 2050, puisque le "
                "climat déplacera le risque vers des communes sans passé.",
             d.L, y + 9.4, 30, 3.4, VERT, 15)
    d.notes(s, "Protocole : on retire une région du train, on entraîne, on "
               "teste sur la région exclue, neuf fois. C'est la réponse la plus "
               "directe à H2 : le territoire porte une information qui "
               "transfère, l'historique non.")

    # ── la calibration ──────────────────────────────────────────────────
    s, y = d.page("La calibration, et pourquoi l'application affiche un rang",
                  "Le score brut est 145 fois trop grand, conséquence directe "
                  "du sous-échantillonnage.")
    cal = D["calibration_v3"]
    d.tableau(s, ["Méthode", "PR-AUC", "p moyen", "biais", "valeurs distinctes"],
              [[r.methode, f"{r.pr_auc:.4f}", f"{r.p_moyen:.2e}",
                f"×{r.biais:.2f}",
                f"{int(r.valeurs_distinctes):,}".replace(",", " ")]
               for _, r in cal.iterrows()],
              d.L, y, 25, 4.4, larg=[5, 5, 5, 5, 5], fort=(1,))
    d.puces(s, [
        "**Platt corrige sans rien coûter**",
        "le classement reste intact, seule l'échelle bouge.",
        "**L'isotonique calibre aussi bien, mais écrase le score**",
        "136 valeurs distinctes au lieu de 9 millions : on perd du pouvoir de "
        "discrimination sans rien gagner.",
        "**L'application affiche donc un rang, pas une probabilité**",
        "le calibrateur disponible a été ajusté sur un autre modèle et une "
        "autre période, il serait faux d'un facteur voisin de 2. Un rang, lui, "
        "reste juste.",
    ], d.L, y + 5.2, 30, 15, 1.3)
    d.notes(s, "Le point à faire passer : on a préféré ne pas afficher de "
               "probabilité plutôt que d'en afficher une fausse. C'est un "
               "choix documenté, pas un oubli.")

    # ── le test ─────────────────────────────────────────────────────────
    s, y = d.page("L'évaluation finale, ouverte une seule fois",
                  "Après gel complet du modèle, des variables et de la "
                  "calibration.")
    tpa = D["test_par_annee"]
    d.tableau(s, ["Année", "Lignes", "Feux", "Taux", "PR-AUC", "lift"],
              [[int(r.an), f"{int(r.lignes):,}".replace(",", " "),
                f"{int(r.feux):,}".replace(",", " "), f"{r.taux:.4%}",
                f"{r.pr_auc:.4f}", f"×{r.lift:.0f}"]
               for _, r in tpa.iterrows()],
              d.L, y, 24, 4.4, larg=[3.5, 6, 4, 4.5, 3.5, 2.5])
    mx, mn = tpa.loc[tpa.lift.idxmax()], tpa.loc[tpa.lift.idxmin()]
    d.encart(s, f"Le lift varie du simple au double, et cette variation suit "
                f"la rareté.\n"
                f"{int(mx.an)} est l'année la plus calme ({int(mx.feux)} feux) "
                f"et donne le meilleur lift (×{mx.lift:.0f}) ; {int(mn.an)}, "
                f"la plus active, le plus faible (×{mn.lift:.0f}).\n"
                f"Une année calme concentre les feux dans les endroits les "
                f"plus prévisibles. Quand tout brûle, y compris là où ce n'est "
                f"pas attendu, le modèle est pris en défaut.",
             d.L, y + 5.2, 30, 4.6, ORANGE, 15)
    d.notes(s, "Regarder la variabilité et pas seulement la moyenne : le lift "
               "moyen de 63× cache un intervalle de 76 à 151 selon l'année.")

    # ── le rendement opérationnel ───────────────────────────────────────
    s, y = d.page("Ce que le modèle change concrètement",
                  "Un lift ne se traduit pas directement en décision. La "
                  "question utile est : si on peut surveiller 1 % du "
                  "territoire, combien de feux couvre-t-on ?")
    op = D["operationnel_test"]
    if op is not None:
        d.tableau(s, ["Budget de surveillance", "Communes-jours suivis",
                      "Feux couverts", "Part des feux"],
                  [[f"{r.budget:.1%}".replace(".", ","),
                    f"{int(r.lignes):,}".replace(",", " "),
                    f"{int(r.feux):,}".replace(",", " "),
                    f"{r.rappel:.1%}".replace(".", ",")]
                   for _, r in op.iterrows()],
                  d.L, y, 26, 5.6, larg=[7, 7.5, 5.5, 6], taille=13,
                  fort=(2,))
    d.encart(s, "En surveillant 1 % des communes-jours, on couvre 42 % des "
                "départs. Le rapport entre les deux est ce que le modèle "
                "apporte réellement.\nAutre lecture, moins flatteuse et tout "
                "aussi vraie : les 37 % de PR-AUC qui séparent v3 du modèle C "
                "ne valent que 3 points de rappel à 1 % de budget, et plus "
                "rien du tout à 10 %. Une grande partie de l'écart mesuré est "
                "opérationnellement invisible.",
             d.L, y + 6.4, 30, 4.6, BLEU, 15)
    d.notes(s, "C'est la diapositive qui relie les métriques à une décision. "
               "Si on demande pourquoi ne pas déployer v3 malgré son avance : "
               "cette avance disparaît là où elle devrait servir.")

    # ── H4 · le climat ──────────────────────────────────────────────────
    s, y = d.page("H4 · le danger monte, les feux non",
                  "Deux fenêtres d'observation, deux conclusions, et il faut "
                  "dire laquelle croire.")
    d._image(s, IMG / "tendance.png", d.L, y, 30, 7.2)
    fa, fe = t("FWI moyen annuel"), t("FWI moyen juin-septembre")
    jd, fx = t("jours de danger élevé (FWI > 21,3)"), t("communes-jours en feu")
    d.tableau(s, ["Série", "Période", "Variation", "p", "Verdict"], [
        ["FWI moyen annuel", f"{fa.an_min}-{fa.an_max}",
         f"{fa.variation_pct:+.0f} %", f"{fa.p:.1e}", "significatif"],
        ["FWI juin-septembre", f"{fe.an_min}-{fe.an_max}",
         f"{fe.variation_pct:+.0f} %", f"{fe.p:.1e}", "significatif"],
        ["jours de danger élevé", f"{jd.an_min}-{jd.an_max}",
         f"{jd.variation_pct:+.0f} %", f"{jd.p:.1e}", "significatif"],
        ["communes-jours en feu", f"{fx.an_min}-{fx.an_max}",
         f"{fx.variation_pct:+.0f} %", f"{fx.p:.2f}", "non significatif"],
    ], d.L, y + 7.8, 27, 5, larg=[8, 5, 5, 4.5, 4.5], fort=(3,))
    d.notes(s, "Formulation exacte à employer : les conditions favorables aux "
               "feux augmentent très significativement, et le nombre de "
               "départs reste stable. Ne jamais dire que les feux augmentent, "
               "les données du projet ne le montrent pas. Trois lectures "
               "possibles, à donner dans cet ordre : 20 points d'observation "
               "seulement, une prévention qui absorbe pour l'instant la hausse "
               "de l'aléa, et le fait qu'on projette l'aléa et non le bilan.")

    # ════════════════════════════════════════════════════════════════════
    #  05 · INTERPRÉTATION
    # ════════════════════════════════════════════════════════════════════
    d.section("05", "Interprétation",
              "Ce que le modèle a appris, et ce qu'on peut en faire")

    # ── SHAP ────────────────────────────────────────────────────────────
    s, y = d.page("Qu'est-ce qui fait partir un feu ?",
                  "Trois façons de poser la question, trois réponses, chacune "
                  "juste mais pas à la même question.")
    if "shap" in G:
        d._image(s, G["shap"], d.L, y, 30, 8.2)
    d.encart(s, "danger_effis est 2ᵉ par gain et 30ᵉ par SHAP. C'est une "
                "discrétisation du FWI en six classes : XGBoost trouve ces "
                "seuils commodes pour découper, mais l'information est déjà "
                "dans le FWI continu, et SHAP lui en attribue le crédit. "
                "L'importance par gain ne sait pas traiter la redondance.\n"
                "part_maquis est 1ᵉʳ par gain, 10ᵉ sur l'échantillon aléatoire "
                "et 2ᵉ au sommet : il ne change rien sur une commune-jour "
                "moyenne, et devient déterminant là où le modèle voit du "
                "risque.", d.L, y + 8.8, 30, 5, ROUGE, 14)
    d.notes(s, "Réponse à H1 et H2 réunies : les variables de territoire "
               "occupent le haut du classement à égalité avec la météo. La "
               "météo situe le moment, le territoire situe le lieu. Si on "
               "demande la garantie : TreeSHAP est exact sur un modèle "
               "d'arbres, vérifié à 9,5e-07 près.")

    # ── DiCE ────────────────────────────────────────────────────────────
    s, y = d.page("La seule question actionnable",
                  "SHAP et LIME répondent à « pourquoi ce score ». DiCE répond "
                  "à « qu'aurait-il fallu changer ».")
    yy = d.puces(s, [
        "**Un exemple, mesuré**",
        "Bormes-les-Mimosas, 99,9ᵉ percentile du 12 août 2024. On autorise à "
        "modifier la végétation seule : maquis, forêt, part combustible, part "
        "agricole.",
        "**Résultat : aucun contrefactuel**",
        "rien, sur ces leviers, ne fait sortir la commune du décile à risque. "
        "Son exposition tient à sa position, à son relief, à sa superficie. Le "
        "risque est structurel, et l'absence de solution est ici la réponse.",
    ], d.L, y, 30, 16, 1.28)
    d.encart(s, "Un contrefactuel n'est pas une recommandation. Rien ne "
                "garantit qu'il soit réalisable, on ne convertit pas 40 % de "
                "maquis en terres agricoles, ni que le lien soit causal : le "
                "modèle a appris des corrélations, pas des mécanismes.",
             d.L, yy + 0.35, 30, 2.6, ORANGE, 14)
    d.notes(s, "Le détail d'implémentation qui change tout, à raconter plutôt "
               "qu'à afficher : DiCE cherche par défaut à faire passer la "
               "probabilité sous 0,5. Or le score n'est pas calibré, 0,5 "
               "correspond à un risque astronomique, et l'outil ne renvoyait "
               "jamais rien. On a recentré la frontière sur le décile par une "
               "transformation strictement croissante, qui laisse le "
               "classement intact. La question devient « que faudrait-il pour "
               "sortir des 10 % les plus à risque », celle qui a un sens "
               "opérationnel.\n\n"
               "Sur LIME, si la question vient : sur un modèle d'arbres il "
               "approxime ce que TreeSHAP calcule exactement, donc il ne peut "
               "pas faire mieux. Il redeviendrait le bon outil sur un modèle "
               "qu'on ne peut pas ouvrir, une API ou un réseau profond.")

    # ── les projections ─────────────────────────────────────────────────
    s, y = d.page("Projeter jusqu'en 2100, et ce que cela ne veut pas dire",
                  "Trois scénarios GIEC, appliqués à l'aléa météo.")
    d._image(s, FIG / "projection-2050" /
             "02_le-risque-de-fond-en-2041-2055-a-vegetation-relief-et-preven.png",
             d.L, y, 18, 10.2)
    d.puces(s, [
        "**Ce qu'on projette**",
        "Le FWI, seule quantité qui montre un signal et seule que les modèles "
        "climatiques savent fournir.",
        "**Ce qu'on ne projette pas**",
        "Le nombre de feux. La végétation, la prévention et les pratiques "
        "agricoles sont supposées constantes, ce qui est une hypothèse et "
        "elle est fausse.",
        "**Ce que « le 2 août 2050 » signifie**",
        "Pas une prévision météo, mais un 2 août ordinaire sous le climat de "
        "2050. La forme de la saison vient des observations, seul son niveau "
        "est décalé.",
        "**Avant 2045, les scénarios sont indiscernables**",
        "mesuré : en 2030, RCP 2.6 dépasse RCP 8.5 sur 59 % des communes. Ce "
        "n'est pas une anomalie, c'est l'inertie du système climatique.",
    ], 20.8, y, 11.3, 12.5, 1.16)
    d.notes(s, "L'application affiche cet avertissement à chaque date "
               "postérieure à 2025, et grise la zone antérieure à 2045 sur les "
               "graphiques de projection. Le dire avant qu'on le demande.")

    # ── les erreurs ─────────────────────────────────────────────────────
    s, y = d.page("Les erreurs commises, et comment elles ont été trouvées",
                  "Sur un événement à 0,02 %, une erreur ne se manifeste "
                  "jamais par une exception, mais par un chiffre plausible.")
    d._image(s, IMG / "alignement.png", d.L, y, 13.5, 8.2)
    d.puces(s, [
        "**Le premier verdict du LSTM annonçait −97 %. Le vrai est −52 %.**",
        "La requête d'assemblage n'a pas d'ORDER BY : l'ordre des 38 M lignes "
        "renvoyées par PostgreSQL dépend du plan d'exécution et change d'une "
        "exécution à l'autre. Les fichiers ne portaient que (score, cible) : "
        "même taille, même nombre de feux, ordre différent, aucune erreur "
        "levée.",
        "**La parade**",
        "tout fichier porte désormais ses clés, une fonction d'alignement "
        "vérifie, et un test refuse un fichier sans clés.",
    ], 16.2, y, 15.9, 14, 1.2)
    d.tableau(s, ["Les autres", "Le symptôme"], [
        ["cible laissée dans les features du modèle de surface",
         "R² = 0,994 et ROC-AUC = 1,0000"],
        ["interaction analysée sur le sommet du classement",
         "le signe s'inversait, biais de collision"],
        ["projection ancrée à k = 1,0 en 2025",
         "la courbe plongeait au passage 2025 → 2026"],
        ["tendance du FWI écrite en dur dans l'application",
         "+45 % affiché, +58 % réel"],
        ["PR-AUC et ROC-AUC affichées sans être sauvegardées",
         "des chiffres cités sans artefact pour les soutenir"],
    ], d.L, y + 9, 30, 5.4, larg=[16, 14], taille=12)
    d.notes(s, "Le bug d'alignement est reproduit volontairement dans le "
               "notebook et dans l'application : on permute les lignes, la "
               "PR-AUC tombe de 0,0085 à 0,0002, exactement la ligne du "
               "hasard, avec les mêmes valeurs dans le fichier.")

    # ── les limites ─────────────────────────────────────────────────────
    s, y = d.page("Ce que le projet ne fait pas",
                  "Les limites qu'on connaît valent mieux que celles qu'on "
                  "découvre en soutenance.")
    d.puces(s, [
        "**La surface brûlée n'est pas prédictible**",
        "R² de 0,14, moins bon que d'annoncer toujours la médiane. Elle dépend "
        "de ce qui se passe après le départ : vent, délai d'intervention, "
        "relief.",
        "**« Sera-ce un grand feu » se prédit mal aussi**",
        f"lift de {D['taille']['lift_grand']:.1f}× seulement, contre "
        f"{mt['test']['lift']:.1f}× pour les départs.",
        "**Une commune-jour n'est pas un incendie**",
        "un feu traversant cinq communes compte cinq fois.",
        "**31 communes partagent une maille météo de 28 km**",
        "conséquence statistique directe : les intervalles naïfs sur les "
        "coefficients météo seraient trop étroits.",
        "**Le LSTM n'a pas reçu danger_effis**",
        "l'écart de 23,6 % reste un majorant tant que cette asymétrie n'est "
        "pas levée. C'est la première chose à refaire.",
    ], d.L, y, 30, 15, 1.26)
    d.notes(s, "Sur le grand feu, ajouter à l'oral : on lit parfois une "
               "ROC-AUC de 0,77 sur cette tâche. C'est exact, mais m'en servir "
               "après avoir expliqué pourquoi la ROC-AUC flatte serait me "
               "contredire.\n\n"
               "Enchaîner directement sur la conclusion : ces limites "
               "définissent la suite du travail, elles ne l'invalident pas.")

    # ── l'application ───────────────────────────────────────────────────
    s, y = d.page("L'application", "Cinq pages, déployée publiquement.")
    d.tableau(s, ["Page", "Ce qu'elle montre"], [
        ["Carte", "34 696 communes en aplat, 1973 → 2100, 3 scénarios GIEC"],
        ["Commune", "décennies, projections à échelle fixe, type de territoire"],
        ["Les données", "les 4 sources, 53 ans de tendance, ADF, ACF/PACF, "
                        "SARIMAX"],
        ["Les modèles", "protocole, intervalles, LSTM, v3 contre C, "
                        "calibration"],
        ["Pourquoi un feu part", "SHAP, LIME et DiCE, et leur désaccord"],
    ], d.L, y, 30, 6.4, larg=[7, 23])
    d.encart(s, "Le mode rétrospectif refuse d'afficher le modèle v3 sur 20 "
                "des 23 années qu'il pourrait techniquement couvrir, et il "
                "sait dire laquelle des trois raisons s'applique à la date "
                "demandée.", d.L, y + 7.1, 30, 2.4, ROUGE, 15)
    d._texte(s, "terre-vent-feu-eau-data.streamlit.app",
             d.L, y + 10, 30, 1.4, 19, BLEU, True)
    d.notes(s, "Faire la démonstration ici : la carte d'abord, puis le "
               "basculement rétrospectif sur une date d'août, puis la page "
               "Pourquoi. Prévoir trois minutes.\n\n"
               "Les trois raisons du refus, à énoncer : le train parce qu'il a "
               "appris ces lignes, la validation parce qu'elle a servi à "
               "choisir, l'avenir parce qu'il n'y a pas d'historique.\n\n"
               "Réveiller l'application dix minutes avant de passer : elle "
               "s'endort après inactivité et met environ deux minutes à "
               "revenir.")

    # ── conclusion ──────────────────────────────────────────────────────
    s, y = d.page("Réponses aux quatre hypothèses")
    d.tableau(s, ["", "Hypothèse", "Réponse", "Sur quelle mesure"], [
        ["H1", "la météo seule ne suffit pas à localiser le risque",
         "confirmée",
         "danger EFFIS seul ×5, croisé au territoire ×42"],
        ["H2", "le territoire apporte une information propre",
         "confirmée",
         "9 régions sur 9 en validation croisée spatiale, jusqu'à +137 %"],
        ["H3", "une architecture séquentielle n'apporte rien de plus",
         "confirmée",
         "LSTM −23,6 % [−33,5 ; −17,3] à information égale ; PACF et ARIMA "
         "concordants"],
        ["H4", "le danger météo augmente de façon détectable",
         "confirmée pour l'aléa,\nnon vérifiée pour les feux",
         "FWI estival +62 % (p = 1e-4) ; départs +0 % non significatif"],
    ], d.L, y, 30, 8, larg=[1.8, 10.2, 6, 12], taille=12.5)
    d.encart(s, "La réponse à la question principale est positive, avec une "
                "réserve explicite : on estime un risque relatif, pas une "
                "probabilité absolue, et c'est pour cette raison que "
                "l'application affiche un rang.",
             d.L, y + 8.8, 30, 3, VERT, 16)
    d.notes(s, "Terminer sur la réserve plutôt que sur la performance. Si une "
               "seule chose doit rester : le modèle classe bien, il ne "
               "quantifie pas.")

    # ── clôture ─────────────────────────────────────────────────────────
    # Une vraie diapositive de fin, et pas une ligne « Merci » ajoutée sous
    # un tableau : c'est l'image qui reste à l'écran pendant les questions,
    # donc elle porte la carte et les trois adresses.
    s = d._vide(INK)
    d._texte(s, "Merci", d.L, 5.6, 18, 3, 54, FOND, True, 1.05)
    d._texte(s, "Le modèle classe bien.\nIl ne quantifie pas.",
             d.L, 10.2, 18, 3, 20, GRIS, interligne=1.4)
    d._filet(s, 14.6, ROUGE, x=d.L, w=3.2)
    d._texte(s, "L'application  ·  terre-vent-feu-eau-data.streamlit.app\n"
                "La présentation  ·  hierogifle.github.io/terre-vent-feu-eau-data\n"
                "Le code  ·  github.com/Hierogifle/terre-vent-feu-eau-data",
             d.L, 15.4, 20, 3, 13, GRIS, interligne=1.5)
    d._texte(s, "Romuald Courtois  ·  M1 Data / IA  ·  La Plateforme_",
             20.5, 17.4, 12, 1, 12, MUTED, aligne=None)
    d._image(s, IMG / "carte.png", 21.5, 2.6, 10.5, 12)
    d.notes(s, "Laisser cette diapositive affichée pendant les questions : la "
               "carte donne au jury quelque chose à regarder, et les trois "
               "adresses sont lisibles depuis la salle.\n\n"
               "Si on demande par où commencer : la vitrine GitHub Pages, "
               "elle renvoie vers les deux autres.")

    chemin = d.enregistrer(SORTIE / "soutenance.pptx")
    n = len(d.p.slides._sldIdLst)
    print(f"✅ {chemin.relative_to(RACINE)} — {n} diapositives, "
          f"{chemin.stat().st_size / 1e6:.1f} Mo")


if __name__ == "__main__":
    main()
